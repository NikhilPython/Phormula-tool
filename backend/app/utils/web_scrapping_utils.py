import os
import trafilatura
from urllib.parse import urljoin, urlparse
from dotenv import load_dotenv
from openai import OpenAI
from playwright.sync_api import sync_playwright
from bs4 import BeautifulSoup
from pptx import Presentation


load_dotenv()
openai_client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

def extract_ppt_text(file_path):

    prs = Presentation(file_path)

    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text


def get_internal_links(base_url):

    # ✅ clean URL (remove tracking params)
    if "?" in base_url:
        base_url = base_url.split("?")[0]

    visited = set()
    to_visit = [base_url]
    pages = []

    domain = urlparse(base_url).netloc

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)

        page = browser.new_page(
            user_agent="Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 Chrome/120 Safari/537.36"
        )

        # ✅ speed boost (block heavy resources)
        page.route("**/*", lambda route: route.abort()
            if route.request.resource_type in ["image", "media", "font"]
            else route.continue_()
        )

        while to_visit and len(visited) < 20:
            url = to_visit.pop(0)

            if url in visited:
                continue

            visited.add(url)

            try:
                # ✅ better loading strategy (fix timeout issue)
                try:
                    page.goto(url, timeout=90000, wait_until="domcontentloaded")
                except:
                    page.goto(url, timeout=90000)

                html = page.content()

                soup = BeautifulSoup(html, "html.parser")

                pages.append(html)

                # find links
                for link in soup.find_all("a", href=True):
                    full_url = urljoin(url, link["href"])

                    # ❗ skip useless links
                    if (
                        "#" in full_url or
                        "mailto:" in full_url or
                        "javascript:" in full_url
                    ):
                        continue

                    # ❗ skip junk pages
                    if any(x in full_url for x in ["privacy", "terms", "login", "signup"]):
                        continue

                    # only same domain
                    if domain in urlparse(full_url).netloc:
                        if full_url not in visited and full_url not in to_visit:
                            to_visit.append(full_url)

            except Exception as e:
                print(f"Error crawling {url}: {e}")

        browser.close()

    return pages

def extract_text(html):

    text = trafilatura.extract(
        html,
        include_comments=False,
        include_tables=False
    )

    if text:
        return text

    # 🔁 fallback to BeautifulSoup
    soup = BeautifulSoup(html, "html.parser")

    for tag in soup(["script","style","nav","footer","header"]):
        tag.extract()

    return soup.get_text(" ", strip=True)
def analyze_business_website(url=None, ppt_path=None):
    # -------- Normalize inputs --------
    url = (url or "").strip()

    if not url and not ppt_path:
        return {
            "error": "Either website or PPT is required"
        }

    # -------- Website --------
    pages = []
    website_text = ""

    if url:
        try:
            pages = get_internal_links(url)

            for html in pages:
                text = extract_text(html)
                if text:
                    website_text += text + "\n"

            website_text = website_text[:15000]

        except Exception as e:
            print(f"Website extraction error: {e}")

    # -------- PPT --------
    ppt_text = ""

    if ppt_path:
        try:
            ppt_text = extract_ppt_text(ppt_path)
            ppt_text = (ppt_text or "")[:5000]
        except Exception as e:
            print(f"PPT extraction error: {e}")

    # -------- Combine --------
    combined_text = (website_text + "\n" + ppt_text).strip()

    if len(combined_text) < 500:
        return {
            "error": "Could not extract sufficient content from website or PPT"
        }

    # -------- Prompt --------
    prompt = f"""
You are the founder of a company.

Using the available content, write a clear and compelling business overview in a natural, human tone.

Guidelines:
- Write in first person (we, our)
- Keep it around 200–250 words
- Make it sound confident, clear, and slightly conversational
- If both website and PPT are available, combine insights from both
- Prefer PPT for vision, positioning, and strategy
- Prefer website for product and operational details
- Avoid repetition
- Avoid buzzwords and generic phrases
- Do NOT mention sources
- Do NOT use bullet points
- Write as a single clean paragraph

Content:
{combined_text}
"""

    response = openai_client.responses.create(
        model="gpt-4.1",
        input=prompt
    )

    ai_output = response.output_text.strip()

    return {
        "overview": ai_output,
        "debug": {
            "website_text_sample": website_text,
            "ppt_text_sample": ppt_text,
            "pages_scraped": len(pages),
            "used_website": bool(url),
            "used_ppt": bool(ppt_path)
        }
    }

   